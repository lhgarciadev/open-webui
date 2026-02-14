import json
import os
import subprocess
import sys
from pathlib import Path


def _backend_dir() -> Path:
    return Path(__file__).resolve().parents[3]


def _extract_json_line(stdout: str) -> dict:
    for line in reversed(stdout.splitlines()):
        line = line.strip()
        if not line:
            continue
        try:
            return json.loads(line)
        except json.JSONDecodeError:
            continue
    raise AssertionError(f"No JSON payload found in subprocess output:\n{stdout}")


def test_presentations_tool_validation_cases():
    env = os.environ.copy()
    env["PYTHONPATH"] = str(_backend_dir())

    script = """
import asyncio
import json
import os
import tempfile
from pathlib import Path

data_dir = tempfile.mkdtemp(prefix="cognitia_presentations_test_")
os.environ["DATA_DIR"] = data_dir

from open_webui.tools.presentations import generate_presentation

async def run():
    slides = [
        {"type": "title", "title": "Demo", "subtitle": "Validacion"},
        {"type": "content", "title": "Puntos", "bullets": ["Uno", "Dos"]},
    ]
    ok = json.loads(await generate_presentation(title="Validacion Tool", slides=slides))
    err = json.loads(await generate_presentation(title="Sin Contenido"))

    path = Path(ok.get("file_path", ""))
    expected_dir = Path(data_dir) / "presentations"
    payload = {
        "case_success_with_slides": {
            "success": ok.get("success") is True,
            "slides_count": ok.get("slides_count") == 2,
            "download_url_matches_filename": ok.get("download_url") == f"/api/v1/files/presentations/{ok.get('filename','')}",
            "file_written_under_data_dir_presentations": path.exists() and path.parent.resolve() == expected_dir.resolve(),
        },
        "case_error_without_slides_or_story_spec": {
            "success": err.get("success") is False,
            "error_matches": err.get("error") == "No slide content provided. Send non-empty 'slides' or 'story_spec'.",
            "hint_present": "story_spec schema" in err.get("hint", ""),
        },
        "case_download_url_and_data_dir_write": {
            "download_url_prefix": str(ok.get("download_url", "")).startswith("/api/v1/files/presentations/"),
            "download_url_exact": ok.get("download_url") == f"/api/v1/files/presentations/{ok.get('filename','')}",
            "file_exists": path.exists(),
            "dir_matches_data_dir_presentations": path.parent.resolve() == expected_dir.resolve(),
        },
    }
    print(json.dumps(payload))

asyncio.run(run())
"""

    result = subprocess.run(
        [sys.executable, "-c", script],
        env=env,
        capture_output=True,
        text=True,
        timeout=90,
    )

    assert result.returncode == 0, result.stderr
    payload = _extract_json_line(result.stdout)

    for case_name, checks in payload.items():
        for check_name, value in checks.items():
            assert value is True, f"{case_name}.{check_name} failed"


def test_story_spec_gamma_rules_and_limits():
    env = os.environ.copy()
    env["PYTHONPATH"] = str(_backend_dir())

    script = """
import asyncio
import json
import os
import tempfile
from pathlib import Path

from pptx import Presentation

data_dir = tempfile.mkdtemp(prefix="cognitia_story_spec_gamma_")
os.environ["DATA_DIR"] = data_dir

from open_webui.tools.presentations import generate_presentation, _build_slides_from_story_spec

story_spec = [
    {
        "type": "insight",
        "title": "Hallazgos",
        "points": [f"Punto {i}" for i in range(1, 12)],
    },
    {
        "type": "metrics",
        "title": "KPIs",
        "metrics": [
            {"value": f"{i * 10}%", "label": f"Metrica {i}"}
            for i in range(1, 7)
        ],
    },
]

slides = _build_slides_from_story_spec("Validacion Gamma", story_spec)
content_slide = next((s for s in slides if s.get("type") == "content"), {})
stats_slide = next((s for s in slides if s.get("type") == "stats"), {})

result = json.loads(asyncio.run(generate_presentation(title="Validacion Gamma", story_spec=story_spec)))
path = Path(result.get("file_path", ""))
expected_dir = Path(data_dir) / "presentations"

render_ok = False
if result.get("success") and path.exists():
    prs = Presentation(str(path))
    render_ok = len(prs.slides) == result.get("slides_count") and len(prs.slides) > 0

payload = {
    "cover_auto_added_when_missing": bool(slides and slides[0].get("type") == "title"),
    "closing_auto_added_when_missing": bool(slides and slides[-1].get("type") == "closing"),
    "bullets_capped_to_8": len(content_slide.get("bullets", [])) == 8,
    "metrics_capped_to_4": len(stats_slide.get("stats", [])) == 4,
    "pptx_generated_and_renderable": render_ok,
    "file_written_under_data_dir_presentations": path.exists() and path.parent.resolve() == expected_dir.resolve(),
}

print(json.dumps(payload))
"""

    result = subprocess.run(
        [sys.executable, "-c", script],
        env=env,
        capture_output=True,
        text=True,
        timeout=90,
    )

    assert result.returncode == 0, result.stderr
    payload = _extract_json_line(result.stdout)

    for check_name, value in payload.items():
        assert value is True, f"{check_name} failed"
