"""
Cognitia Slides - Presentation generation tool.

Enables AI agents to create professional PowerPoint presentations
with brand-consistent templates and styling.
"""

import json
import logging
import os
import uuid
from datetime import datetime
from typing import Optional, List, Dict, Any
from pathlib import Path

from fastapi import Request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

log = logging.getLogger(__name__)

# =============================================================================
# BRAND CONFIGURATION - Cognitia
# =============================================================================

BRAND_COLORS = {
    "primary": RGBColor(59, 130, 246),      # #3b82f6 - Blue 500
    "primary_dark": RGBColor(30, 64, 175),   # #1e40af - Blue 800
    "surface": RGBColor(15, 23, 42),         # #0f172a - Slate 900
    "surface_light": RGBColor(30, 41, 59),   # #1e293b - Slate 800
    "text_primary": RGBColor(248, 250, 252), # #f8fafc - Slate 50
    "text_secondary": RGBColor(148, 163, 184), # #94a3b8 - Slate 400
    "accent_success": RGBColor(34, 197, 94),  # #22c55e - Green 500
    "accent_warning": RGBColor(245, 158, 11), # #f59e0b - Amber 500
    "accent_error": RGBColor(239, 68, 68),    # #ef4444 - Red 500
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
}

AVAILABLE_ICONS = [
    {"name": "chart", "description": "Bar chart icon for data/analytics"},
    {"name": "users", "description": "People icon for team/users"},
    {"name": "rocket", "description": "Rocket icon for launch/growth"},
    {"name": "target", "description": "Target icon for goals/objectives"},
    {"name": "lightbulb", "description": "Lightbulb icon for ideas/innovation"},
    {"name": "shield", "description": "Shield icon for security/protection"},
    {"name": "globe", "description": "Globe icon for global/international"},
    {"name": "clock", "description": "Clock icon for time/schedule"},
    {"name": "check", "description": "Checkmark icon for completion/success"},
    {"name": "star", "description": "Star icon for highlights/favorites"},
    {"name": "heart", "description": "Heart icon for favorites/health"},
    {"name": "gear", "description": "Gear icon for settings/technical"},
    {"name": "database", "description": "Database icon for data/storage"},
    {"name": "cloud", "description": "Cloud icon for cloud services"},
    {"name": "lock", "description": "Lock icon for security/privacy"},
]

SLIDE_TEMPLATES = [
    "title",           # Title slide with subtitle
    "content",         # Title + bullet points
    "two_column",      # Two column layout
    "section",         # Section divider
    "image",           # Title + image placeholder
    "comparison",      # Side by side comparison
    "quote",           # Quote/testimonial
    "stats",           # Statistics/metrics display
    "timeline",        # Timeline/roadmap
    "closing",         # Thank you/closing slide
]

STORY_SPEC_TYPES = [
    "cover",
    "section",
    "insight",
    "metrics",
    "comparison",
    "timeline",
    "quote",
    "cta",
]

# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def _get_presentations_dir() -> Path:
    """Get the presentations output directory."""
    base_dir = Path(os.environ.get("DATA_DIR", "/app/backend/data"))
    presentations_dir = base_dir / "presentations"
    presentations_dir.mkdir(parents=True, exist_ok=True)
    return presentations_dir


def _apply_brand_background(slide, dark: bool = True):
    """Apply brand background color to slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_COLORS["surface"] if dark else BRAND_COLORS["white"]


def _add_title_shape(slide, text: str, top: float = 0.5, font_size: int = 44, dark: bool = True):
    """Add a styled title to the slide."""
    left = Inches(0.5)
    top = Inches(top)
    width = Inches(9)
    height = Inches(1)

    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["text_primary"] if dark else BRAND_COLORS["surface"]
    p.alignment = PP_ALIGN.LEFT


def _add_subtitle_shape(slide, text: str, top: float = 1.5, dark: bool = True):
    """Add a styled subtitle to the slide."""
    left = Inches(0.5)
    top = Inches(top)
    width = Inches(9)
    height = Inches(0.75)

    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.color.rgb = BRAND_COLORS["text_secondary"] if dark else BRAND_COLORS["surface_light"]
    p.alignment = PP_ALIGN.LEFT


def _add_bullets(slide, items: List[str], top: float = 2.5, dark: bool = True):
    """Add bullet points to the slide."""
    left = Inches(0.5)
    top = Inches(top)
    width = Inches(9)
    height = Inches(4)

    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = BRAND_COLORS["text_primary"] if dark else BRAND_COLORS["surface"]
        p.space_before = Pt(12)
        p.space_after = Pt(6)


def _add_brand_footer(slide, dark: bool = True):
    """Add Cognitia brand footer."""
    left = Inches(0.5)
    top = Inches(7)
    width = Inches(2)
    height = Inches(0.3)

    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Cognitia"
    p.font.size = Pt(10)
    p.font.color.rgb = BRAND_COLORS["primary"]
    p.font.bold = True


def _add_accent_bar(slide):
    """Add brand accent bar at top of slide."""
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(0.1)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_COLORS["primary"]
    shape.line.fill.background()


# =============================================================================
# SLIDE CREATION FUNCTIONS
# =============================================================================

def _create_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    """Create a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    _apply_brand_background(slide, dark=True)

    # Centered title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)

    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["text_primary"]
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        # Subtitle
        left = Inches(0.5)
        top = Inches(4.2)
        width = Inches(9)
        height = Inches(1)

        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = BRAND_COLORS["text_secondary"]
        p.alignment = PP_ALIGN.CENTER

    _add_accent_bar(slide)
    return slide


def _create_content_slide(prs: Presentation, title: str, bullets: List[str]):
    """Create a content slide with title and bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    _apply_brand_background(slide, dark=True)
    _add_accent_bar(slide)

    _add_title_shape(slide, title, top=0.5)
    _add_bullets(slide, bullets, top=1.8)
    _add_brand_footer(slide)

    return slide


def _create_two_column_slide(prs: Presentation, title: str, left_items: List[str], right_items: List[str], left_title: str = "", right_title: str = ""):
    """Create a two-column slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    _apply_brand_background(slide, dark=True)
    _add_accent_bar(slide)

    _add_title_shape(slide, title, top=0.5)

    # Left column title
    if left_title:
        left = Inches(0.5)
        top = Inches(1.6)
        width = Inches(4.2)
        height = Inches(0.5)
        shape = slide.shapes.add_textbox(left, top, width, height)
        p = shape.text_frame.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BRAND_COLORS["primary"]

    # Left column items
    left = Inches(0.5)
    top = Inches(2.2)
    width = Inches(4.2)
    height = Inches(4)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    for i, item in enumerate(left_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = BRAND_COLORS["text_primary"]
        p.space_before = Pt(8)

    # Right column title
    if right_title:
        left = Inches(5.2)
        top = Inches(1.6)
        width = Inches(4.2)
        height = Inches(0.5)
        shape = slide.shapes.add_textbox(left, top, width, height)
        p = shape.text_frame.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BRAND_COLORS["primary"]

    # Right column items
    left = Inches(5.2)
    top = Inches(2.2)
    width = Inches(4.2)
    height = Inches(4)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    for i, item in enumerate(right_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = BRAND_COLORS["text_primary"]
        p.space_before = Pt(8)

    _add_brand_footer(slide)
    return slide


def _create_section_slide(prs: Presentation, title: str, subtitle: str = ""):
    """Create a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Use primary color as background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_COLORS["primary"]

    # Centered title
    left = Inches(0.5)
    top = Inches(3)
    width = Inches(9)
    height = Inches(1.5)

    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        left = Inches(0.5)
        top = Inches(4.5)
        width = Inches(9)
        height = Inches(0.75)

        shape = slide.shapes.add_textbox(left, top, width, height)
        p = shape.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = BRAND_COLORS["white"]
        p.alignment = PP_ALIGN.CENTER

    return slide


def _create_stats_slide(prs: Presentation, title: str, stats: List[dict]):
    """Create a statistics slide. Stats format: [{"value": "50%", "label": "Growth"}]"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    _apply_brand_background(slide, dark=True)
    _add_accent_bar(slide)

    _add_title_shape(slide, title, top=0.5)

    # Calculate positions for stats (max 4)
    num_stats = min(len(stats), 4)
    total_width = 9
    stat_width = total_width / num_stats

    for i, stat in enumerate(stats[:4]):
        left = Inches(0.5 + i * stat_width)
        top = Inches(2.5)
        width = Inches(stat_width - 0.2)
        height = Inches(3)

        # Value
        shape = slide.shapes.add_textbox(left, top, width, Inches(1.5))
        p = shape.text_frame.paragraphs[0]
        p.text = str(stat.get("value", ""))
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = BRAND_COLORS["primary"]
        p.alignment = PP_ALIGN.CENTER

        # Label
        shape = slide.shapes.add_textbox(left, Inches(4), width, Inches(1))
        p = shape.text_frame.paragraphs[0]
        p.text = str(stat.get("label", ""))
        p.font.size = Pt(18)
        p.font.color.rgb = BRAND_COLORS["text_secondary"]
        p.alignment = PP_ALIGN.CENTER

    _add_brand_footer(slide)
    return slide


def _create_quote_slide(prs: Presentation, quote: str, author: str = "", role: str = ""):
    """Create a quote/testimonial slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    _apply_brand_background(slide, dark=True)

    # Quote
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3)

    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(32)
    p.font.italic = True
    p.font.color.rgb = BRAND_COLORS["text_primary"]
    p.alignment = PP_ALIGN.CENTER

    # Author
    if author:
        left = Inches(1)
        top = Inches(5.5)
        width = Inches(8)
        height = Inches(0.75)

        shape = slide.shapes.add_textbox(left, top, width, height)
        p = shape.text_frame.paragraphs[0]
        author_text = f"— {author}"
        if role:
            author_text += f", {role}"
        p.text = author_text
        p.font.size = Pt(20)
        p.font.color.rgb = BRAND_COLORS["primary"]
        p.alignment = PP_ALIGN.CENTER

    _add_brand_footer(slide)
    return slide


def _create_closing_slide(prs: Presentation, title: str = "Thank You", subtitle: str = "", contact: str = ""):
    """Create a closing/thank you slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    _apply_brand_background(slide, dark=True)

    # Title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)

    shape = slide.shapes.add_textbox(left, top, width, height)
    p = shape.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["text_primary"]
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        left = Inches(0.5)
        top = Inches(4.2)
        width = Inches(9)
        height = Inches(0.75)

        shape = slide.shapes.add_textbox(left, top, width, height)
        p = shape.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = BRAND_COLORS["text_secondary"]
        p.alignment = PP_ALIGN.CENTER

    if contact:
        left = Inches(0.5)
        top = Inches(5.5)
        width = Inches(9)
        height = Inches(0.5)

        shape = slide.shapes.add_textbox(left, top, width, height)
        p = shape.text_frame.paragraphs[0]
        p.text = contact
        p.font.size = Pt(18)
        p.font.color.rgb = BRAND_COLORS["primary"]
        p.alignment = PP_ALIGN.CENTER

    _add_accent_bar(slide)
    return slide


def _as_list(value: Any) -> List[Any]:
    """Normalize unknown values into list form."""
    if value is None:
        return []
    if isinstance(value, list):
        return value
    return [value]


def _parse_json_like_list(value: Any) -> List[Dict[str, Any]]:
    """
    Normalize tool inputs that may arrive as list, dict or JSON-encoded string.
    Returns an empty list for unsupported/invalid formats.
    """
    if value is None:
        return []

    parsed_value = value
    if isinstance(value, str):
        raw = value.strip()
        if not raw:
            return []
        try:
            parsed_value = json.loads(raw)
        except Exception:
            return []

    if isinstance(parsed_value, dict):
        parsed_value = [parsed_value]

    if not isinstance(parsed_value, list):
        return []

    return [item for item in parsed_value if isinstance(item, dict)]


def _slide_has_renderable_content(slide_def: Dict[str, Any]) -> bool:
    """Return True when slide has meaningful payload for the selected type."""
    slide_type = str(slide_def.get("type", "content")).lower()

    if slide_type == "title":
        return bool(str(slide_def.get("title", "")).strip() or str(slide_def.get("subtitle", "")).strip())

    if slide_type in {"content", "timeline"}:
        bullets = [
            str(item).strip()
            for item in _as_list(slide_def.get("bullets", []))
            if str(item).strip()
        ]
        return bool(str(slide_def.get("title", "")).strip() or bullets)

    if slide_type == "two_column":
        left_items = [
            str(item).strip()
            for item in _as_list(slide_def.get("left_items", []))
            if str(item).strip()
        ]
        right_items = [
            str(item).strip()
            for item in _as_list(slide_def.get("right_items", []))
            if str(item).strip()
        ]
        return bool(
            str(slide_def.get("title", "")).strip()
            or str(slide_def.get("left_title", "")).strip()
            or str(slide_def.get("right_title", "")).strip()
            or left_items
            or right_items
        )

    if slide_type == "section":
        return bool(str(slide_def.get("title", "")).strip() or str(slide_def.get("subtitle", "")).strip())

    if slide_type == "stats":
        stats = []
        for stat in _as_list(slide_def.get("stats", [])):
            if isinstance(stat, dict):
                value = str(stat.get("value", "")).strip()
                label = str(stat.get("label", "")).strip()
                if value or label:
                    stats.append(stat)
        return bool(str(slide_def.get("title", "")).strip() or stats)

    if slide_type == "quote":
        return bool(
            str(slide_def.get("quote", "")).strip()
            or str(slide_def.get("author", "")).strip()
            or str(slide_def.get("role", "")).strip()
        )

    if slide_type == "closing":
        return bool(
            str(slide_def.get("title", "")).strip()
            or str(slide_def.get("subtitle", "")).strip()
            or str(slide_def.get("contact", "")).strip()
        )

    # Unknown types are treated as content-like slides.
    bullets = [
        str(item).strip()
        for item in _as_list(slide_def.get("bullets", []))
        if str(item).strip()
    ]
    return bool(str(slide_def.get("title", "")).strip() or bullets)


def _auto_generate_slides_from_title(title: str) -> List[dict]:
    """
    Auto-generate a basic presentation structure when no slides/story_spec provided.
    Uses the title to create a sensible default presentation.
    """
    # Parse title to extract potential topic keywords
    title_words = title.lower().split()

    # Default slide structure based on common presentation patterns
    slides = [
        {
            "type": "title",
            "title": title,
            "subtitle": "Presentación generada por Cognitia AI",
        },
        {
            "type": "content",
            "title": "Introducción",
            "bullets": [
                f"Contexto y relevancia de: {title}",
                "Objetivos principales de esta presentación",
                "Estructura del contenido a cubrir",
            ],
        },
        {
            "type": "content",
            "title": "Puntos Clave",
            "bullets": [
                "Punto principal 1 - Personalizar según el tema",
                "Punto principal 2 - Agregar datos específicos",
                "Punto principal 3 - Incluir ejemplos relevantes",
                "Punto principal 4 - Considerar el contexto local",
            ],
        },
        {
            "type": "stats",
            "title": "Datos Relevantes",
            "stats": [
                {"value": "---", "label": "Métrica 1"},
                {"value": "---", "label": "Métrica 2"},
                {"value": "---", "label": "Métrica 3"},
            ],
        },
        {
            "type": "content",
            "title": "Recomendaciones",
            "bullets": [
                "Acción recomendada 1",
                "Acción recomendada 2",
                "Próximos pasos a seguir",
            ],
        },
        {
            "type": "closing",
            "title": "¿Preguntas?",
            "subtitle": f"Gracias por su atención - {title}",
            "contact": "Generado con Cognitia AI",
        },
    ]

    return slides


def _build_slides_from_story_spec(
    title: str,
    story_spec: List[Dict[str, Any]],
) -> List[dict]:
    """
    Convert a Gamma-style narrative spec into internal slide definitions.
    """
    slides: List[dict] = []

    if not story_spec:
        return slides

    first_type = str(story_spec[0].get("type", "")).lower()
    if first_type != "cover":
        slides.append(
            {
                "type": "title",
                "title": title,
                "subtitle": "Executive presentation generated from story spec",
            }
        )

    for block in story_spec:
        block_type = str(block.get("type", "insight")).lower()
        block_title = str(block.get("title", "")).strip() or "Untitled"

        if block_type == "cover":
            slides.append(
                {
                    "type": "title",
                    "title": block_title or title,
                    "subtitle": str(block.get("subtitle", "")).strip(),
                }
            )
            continue

        if block_type == "section":
            slides.append(
                {
                    "type": "section",
                    "title": block_title,
                    "subtitle": str(block.get("subtitle", "")).strip(),
                }
            )
            continue

        if block_type == "metrics":
            stats = []
            for metric in _as_list(block.get("metrics", [])):
                if isinstance(metric, dict):
                    stats.append(
                        {
                            "value": str(metric.get("value", "")).strip(),
                            "label": str(metric.get("label", "")).strip(),
                        }
                    )
            if stats:
                slides.append({"type": "stats", "title": block_title, "stats": stats[:4]})
                continue

        if block_type == "comparison":
            slides.append(
                {
                    "type": "two_column",
                    "title": block_title,
                    "left_title": str(block.get("left_title", "Option A")).strip(),
                    "right_title": str(block.get("right_title", "Option B")).strip(),
                    "left_items": [str(i).strip() for i in _as_list(block.get("left_items", [])) if str(i).strip()],
                    "right_items": [str(i).strip() for i in _as_list(block.get("right_items", [])) if str(i).strip()],
                }
            )
            continue

        if block_type == "timeline":
            milestones = []
            for milestone in _as_list(block.get("milestones", [])):
                if isinstance(milestone, dict):
                    label = str(milestone.get("label", "")).strip()
                    detail = str(milestone.get("detail", "")).strip()
                    if label and detail:
                        milestones.append(f"{label}: {detail}")
                    elif label:
                        milestones.append(label)
            slides.append(
                {
                    "type": "content",
                    "title": block_title,
                    "bullets": milestones[:8],
                }
            )
            continue

        if block_type == "quote":
            slides.append(
                {
                    "type": "quote",
                    "quote": str(block.get("quote", "")).strip(),
                    "author": str(block.get("author", "")).strip(),
                    "role": str(block.get("role", "")).strip(),
                }
            )
            continue

        if block_type == "cta":
            slides.append(
                {
                    "type": "closing",
                    "title": block_title or "Next Steps",
                    "subtitle": str(block.get("subtitle", "")).strip(),
                    "contact": str(block.get("contact", "")).strip(),
                }
            )
            continue

        # Default block type: insight
        bullets = [
            str(i).strip()
            for i in _as_list(block.get("points", []))
            if str(i).strip()
        ]
        if not bullets:
            for k in ["insight", "evidence", "implication", "next_step"]:
                val = str(block.get(k, "")).strip()
                if val:
                    bullets.append(val)

        slides.append(
            {
                "type": "content",
                "title": block_title,
                "bullets": bullets[:8],
            }
        )

    if slides and slides[-1].get("type") != "closing":
        slides.append(
            {
                "type": "closing",
                "title": "Next Steps",
                "subtitle": "Questions and alignment",
                "contact": "",
            }
        )

    return slides


# =============================================================================
# PUBLIC TOOL FUNCTIONS
# =============================================================================

async def get_available_templates(
    __request__: Request = None,
    __user__: dict = None,
) -> str:
    """
    Get a list of available presentation slide templates.

    :return: JSON with available templates and their descriptions
    """
    templates = [
        {"type": "title", "description": "Title slide with main heading and subtitle"},
        {"type": "content", "description": "Insight slide with key points"},
        {"type": "two_column", "description": "Comparison slide (Option A vs Option B)"},
        {"type": "section", "description": "Section divider with accent color background"},
        {"type": "stats", "description": "Metrics slide with prominent KPI cards"},
        {"type": "quote", "description": "Quote or testimonial slide"},
        {"type": "closing", "description": "Call-to-action / closing slide"},
    ]

    return json.dumps({
        "templates": templates,
        "brand": "Cognitia",
        "color_scheme": "Blue corporate theme",
        "story_spec_types": STORY_SPEC_TYPES,
    }, ensure_ascii=False)


async def get_available_icons(
    __request__: Request = None,
    __user__: dict = None,
) -> str:
    """
    Get a list of available icons that can be used in presentations.

    :return: JSON with available icon names and descriptions
    """
    return json.dumps({
        "icons": AVAILABLE_ICONS,
        "note": "Icons are represented as styled shapes in presentations"
    }, ensure_ascii=False)


async def get_story_spec_template(
    __request__: Request = None,
    __user__: dict = None,
) -> str:
    """
    Return a Gamma-style story-spec schema for high-quality narrative slides.
    """
    return json.dumps(
        {
            "story_spec_types": STORY_SPEC_TYPES,
            "guidelines": [
                "Start with cover, then section/insight/metrics blocks, end with cta.",
                "Each block should communicate one idea with concrete evidence.",
                "Prefer short, specific points (max 8 bullets per slide).",
                "Use metrics and comparison blocks to improve decision clarity.",
            ],
            "example": [
                {
                    "type": "cover",
                    "title": "Q1 Product Strategy",
                    "subtitle": "Growth and efficiency priorities",
                },
                {
                    "type": "insight",
                    "title": "Primary challenge",
                    "points": [
                        "Activation rate is below target in SMB segment",
                        "Onboarding friction concentrates in first 24 hours",
                    ],
                },
                {
                    "type": "metrics",
                    "title": "Current KPI baseline",
                    "metrics": [
                        {"value": "42%", "label": "Activation"},
                        {"value": "18%", "label": "Week-4 Retention"},
                        {"value": "$23", "label": "CAC"},
                    ],
                },
                {
                    "type": "comparison",
                    "title": "Strategic options",
                    "left_title": "Optimize onboarding",
                    "left_items": ["Fast impact", "Low implementation risk"],
                    "right_title": "Expand acquisition",
                    "right_items": ["Higher upside", "Higher CAC volatility"],
                },
                {
                    "type": "cta",
                    "title": "Decision requested",
                    "subtitle": "Approve onboarding-first plan for next 2 sprints",
                },
            ],
        },
        ensure_ascii=False,
    )


async def generate_presentation(
    title: str,
    slides: Optional[List[dict] | Dict[str, Any] | str] = None,
    story_spec: Optional[List[dict] | Dict[str, Any] | str] = None,
    author: str = "Cognitia AI",
    __request__: Request = None,
    __user__: dict = None,
) -> str:
    """
    Generate a professional PowerPoint presentation with Cognitia branding.

    IMPORTANT: You MUST provide the 'slides' parameter with actual content for best results.
    If only 'title' is provided, a basic template presentation will be auto-generated.

    :param title: REQUIRED. The presentation title (e.g., "Inteligencia Artificial en Colombia")
    :param slides: RECOMMENDED. List of slide definitions. Each slide needs:
        - type: "title" | "content" | "two_column" | "section" | "stats" | "quote" | "closing"
        - title: The slide heading
        - For "content" type: "bullets" = ["Point 1", "Point 2", "Point 3"]
        - For "stats" type: "stats" = [{"value": "50%", "label": "Growth"}]
        - For "two_column" type: "left_items" and "right_items" = lists of strings
        - For "quote" type: "quote", "author", "role"
        - For "closing" type: "subtitle", "contact"
    :param story_spec: Alternative Gamma-style narrative blocks (auto-converted to slides)
    :param author: Author name for metadata (default: "Cognitia AI")
    :return: JSON with file path, download_url, and presentation details

    CORRECT USAGE EXAMPLE - Always include slides with content:
    generate_presentation(
        title="Inteligencia Artificial en Colombia",
        slides=[
            {"type": "title", "title": "IA en Colombia", "subtitle": "Estado actual y perspectivas 2024"},
            {"type": "content", "title": "Adopción de IA", "bullets": ["45% de empresas usan IA", "Sector financiero lidera", "Crecimiento anual del 30%"]},
            {"type": "stats", "title": "Métricas Clave", "stats": [{"value": "45%", "label": "Adopción"}, {"value": "30%", "label": "Crecimiento"}]},
            {"type": "content", "title": "Desafíos", "bullets": ["Falta de talento", "Infraestructura limitada", "Regulación pendiente"]},
            {"type": "closing", "title": "Conclusiones", "subtitle": "La IA transformará la economía colombiana"}
        ]
    )

    MINIMAL USAGE (auto-generates template slides):
    generate_presentation(title="Inteligencia Artificial en Colombia")
    """
    try:
        normalized_slides = _parse_json_like_list(slides)
        normalized_story_spec = _parse_json_like_list(story_spec)

        if not normalized_slides and normalized_story_spec:
            normalized_slides = _build_slides_from_story_spec(title, normalized_story_spec)

        valid_slides = [
            slide_def
            for slide_def in normalized_slides
            if _slide_has_renderable_content(slide_def)
        ]

        # Auto-generate slides if none provided but we have a title
        if not valid_slides and title.strip():
            log.info(f"Auto-generating slides for presentation: {title}")
            auto_slides = _auto_generate_slides_from_title(title)
            valid_slides = [
                slide_def
                for slide_def in auto_slides
                if _slide_has_renderable_content(slide_def)
            ]

        if not valid_slides:
            return json.dumps(
                {
                    "success": False,
                    "error": "No slide content provided and title is empty.",
                    "hint": "Provide a title and optionally slides/story_spec. If only title is given, slides will be auto-generated.",
                },
                ensure_ascii=False,
            )

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Set metadata
        prs.core_properties.title = title
        prs.core_properties.author = author
        prs.core_properties.comments = "Generated by Cognitia AI"

        # Process each slide
        for slide_def in valid_slides:
            slide_type = slide_def.get("type", "content")

            if slide_type == "title":
                _create_title_slide(
                    prs,
                    slide_def.get("title", title),
                    slide_def.get("subtitle", "")
                )

            elif slide_type == "content":
                _create_content_slide(
                    prs,
                    slide_def.get("title", ""),
                    slide_def.get("bullets", [])
                )

            elif slide_type == "two_column":
                _create_two_column_slide(
                    prs,
                    slide_def.get("title", ""),
                    slide_def.get("left_items", []),
                    slide_def.get("right_items", []),
                    slide_def.get("left_title", ""),
                    slide_def.get("right_title", "")
                )

            elif slide_type == "section":
                _create_section_slide(
                    prs,
                    slide_def.get("title", ""),
                    slide_def.get("subtitle", "")
                )

            elif slide_type == "stats":
                _create_stats_slide(
                    prs,
                    slide_def.get("title", ""),
                    slide_def.get("stats", [])
                )

            elif slide_type == "quote":
                _create_quote_slide(
                    prs,
                    slide_def.get("quote", ""),
                    slide_def.get("author", ""),
                    slide_def.get("role", "")
                )

            elif slide_type == "closing":
                _create_closing_slide(
                    prs,
                    slide_def.get("title", "Thank You"),
                    slide_def.get("subtitle", ""),
                    slide_def.get("contact", "")
                )

        # Generate filename and save
        safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()
        safe_title = safe_title.replace(' ', '_')[:50]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{safe_title}_{timestamp}.pptx"

        output_dir = _get_presentations_dir()
        filepath = output_dir / filename
        prs.save(str(filepath))

        # Build absolute download URL to prevent model from adding sandbox: prefix
        download_path = f"/api/v1/files/presentations/{filename}"

        # Try to get base URL from request or app config
        base_url = ""
        if __request__:
            try:
                # Get base URL from request
                base_url = str(__request__.base_url).rstrip('/')
            except Exception:
                pass

            # Fallback to app config if available
            if not base_url or base_url == "http://testserver":
                try:
                    webui_url = __request__.app.state.config.WEBUI_URL
                    if webui_url:
                        base_url = webui_url.rstrip('/')
                except Exception:
                    pass

        # Build full URL - use absolute URL to prevent sandbox: prefix
        full_download_url = f"{base_url}{download_path}" if base_url else download_path

        return json.dumps({
            "success": True,
            "file_path": str(filepath),
            "filename": filename,
            "slides_count": len(valid_slides),
            "download_url": full_download_url,
            "message": f"Presentation '{title}' created successfully with {len(valid_slides)} slides. Download link: {full_download_url}"
        }, ensure_ascii=False)

    except Exception as e:
        log.exception(f"generate_presentation error: {e}")
        return json.dumps({
            "success": False,
            "error": str(e)
        }, ensure_ascii=False)
